"""CARD-031: PPTXParser — PowerPoint presentation parser."""

from __future__ import annotations

from pathlib import Path

from docforge.config import ParseConfig
from docforge.exceptions import ParseError
from docforge.models import ContentBlock, ContentType
from docforge.parsers import BaseParser, register_parser

try:
    from pptx import Presentation as PptxPresentation
    from pptx.table import Table as PptxTable
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False


class PPTXParser(BaseParser):
    """Parser for PPTX (PowerPoint) presentations.

    Uses python-pptx to traverse slides. Each slide gets a heading,
    shapes are processed as text/tables/pictures.
    """

    def can_parse(self, file_path: Path, file_type: str) -> bool:
        return file_type == "pptx"

    def parse(self, file_path: Path, config: ParseConfig) -> list[ContentBlock]:
        if not _PPTX_AVAILABLE:
            raise ParseError(
                "python-pptx is required to parse PPTX files. "
                "Install with: pip install docforge[pptx]",
                str(file_path),
            )

        prs = PptxPresentation(str(file_path))
        blocks: list[ContentBlock] = []
        reading_order = 0
        image_seq = 0
        image_dir: Path | None = None

        if config.extract_images:
            image_dir = self._resolve_image_dir(file_path, config)
            image_dir.mkdir(parents=True, exist_ok=True)

        for slide_num, slide in enumerate(prs.slides, start=1):
            # Slide heading
            blocks.append(
                ContentBlock(
                    type=ContentType.TEXT,
                    content=f"## Slide {slide_num}",
                    page=slide_num,
                    reading_order=reading_order,
                    metadata={"is_heading": True, "heading_level": 2, "slide_num": slide_num},
                )
            )
            reading_order += 1

            # Process shapes in order
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        blocks.append(
                            ContentBlock(
                                type=ContentType.TEXT,
                                content=text,
                                page=slide_num,
                                reading_order=reading_order,
                            )
                        )
                        reading_order += 1

                if shape.has_table:
                    table = shape.table
                    markdown_table = self._table_to_markdown(table)
                    if markdown_table:
                        blocks.append(
                            ContentBlock(
                                type=ContentType.TABLE,
                                content=markdown_table,
                                page=slide_num,
                                reading_order=reading_order,
                            )
                        )
                        reading_order += 1

                # Picture extraction
                if config.extract_images and image_dir is not None:
                    image_block = self._extract_picture(
                        shape, slide_num, reading_order, image_dir, image_seq
                    )
                    if image_block is not None:
                        blocks.append(image_block)
                        reading_order += 1
                        image_seq += 1

        return blocks

    def _table_to_markdown(self, table: PptxTable) -> str:
        """Convert a python-pptx Table to a Markdown table string."""
        rows_data: list[list[str]] = []
        for row in table.rows:
            row_cells: list[str] = []
            for cell in row.cells:
                row_cells.append(cell.text.strip())
            rows_data.append(row_cells)

        if not rows_data:
            return ""

        # Normalize column count
        max_cols = max(len(row) for row in rows_data)
        for row in rows_data:
            while len(row) < max_cols:
                row.append("")

        lines: list[str] = []
        for row_idx, row_cells in enumerate(rows_data):
            line = "| " + " | ".join(self._escape_cell(cell) for cell in row_cells) + " |"
            lines.append(line)
            if row_idx == 0:
                sep = "| " + " | ".join("---" for _ in row_cells) + " |"
                lines.append(sep)

        return "\n".join(lines)

    def _extract_picture(
        self,
        shape,
        page: int,
        reading_order: int,
        image_dir: Path,
        image_seq: int,
    ) -> ContentBlock | None:
        """Extract a picture shape to an image file."""
        try:
            # Check if shape is a picture
            if shape.shape_type is None:
                return None
            # MSO_SHAPE_TYPE.PICTURE = 13
            if int(shape.shape_type) != 13:  # PICTURE
                return None
        except Exception:
            return None

        try:
            image_bytes = shape.image.blob
            ext = shape.image.content_type.split("/")[-1] if shape.image.content_type else "png"
            if ext.lower() in ("jpeg",):
                ext = "jpg"
            elif ext.lower() not in ("png", "jpg", "gif", "bmp", "svg+xml"):
                ext = "png"
            if ext == "svg+xml":
                ext = "svg"

            filename = f"pptx_img_{image_seq}.{ext}"
            filepath = image_dir / filename
            filepath.write_bytes(image_bytes)

            rel_path = f"{image_dir.name}/{filename}"
            return ContentBlock(
                type=ContentType.IMAGE,
                content=rel_path,
                page=page,
                reading_order=reading_order,
            )
        except Exception:
            return None

    def _resolve_image_dir(self, file_path: Path, config: ParseConfig) -> Path:
        """Resolve the image output directory."""
        if config.image_output_dir:
            return Path(config.image_output_dir)
        stem = file_path.stem
        return file_path.parent / f"{stem}_images"

    @staticmethod
    def _escape_cell(text: str) -> str:
        """Escape pipe characters in table cell text."""
        return text.replace("|", "\\|")


register_parser("pptx", PPTXParser)
