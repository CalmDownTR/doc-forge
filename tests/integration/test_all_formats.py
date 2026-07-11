"""Tests for CARD-033: Unified Routing Verification.

Integration tests verifying all format types work through parse().
"""

from __future__ import annotations

from pathlib import Path

import fitz
import pytest
from docx import Document as DocxDocument
from openpyxl import Workbook
from PIL import Image

from docforge.api import parse
from docforge.exceptions import FileNotSupportedError
from docforge.models import ParseResult
from docforge.ocr import _BACKEND_REGISTRY, OCRBackend, OCRTextResult, register_backend
from docforge.utils.file_utils import detect_file_type


class _DummyOCRBackend(OCRBackend):
    """A dummy OCR backend for integration tests."""

    def initialize(self, languages: list[str]) -> None:
        pass

    def recognize(self, image) -> list[OCRTextResult]:
        return [
            OCRTextResult(
                text="Test OCR result",
                bbox=[[10, 10], [100, 10], [100, 30], [10, 30]],
            ),
        ]

    def is_available(self) -> bool:
        return True


# ============================================================
# Test fixtures: create test files for each format
# ============================================================


def _create_txt_file(tmp_path: Path) -> Path:
    filepath = tmp_path / "test_doc.txt"
    filepath.write_text("Hello World\nThis is a test document.", encoding="utf-8")
    return filepath


def _create_md_file(tmp_path: Path) -> Path:
    filepath = tmp_path / "test_doc.md"
    filepath.write_text("# Title\n\nSome **markdown** content.", encoding="utf-8")
    return filepath


def _create_pdf_file(tmp_path: Path) -> Path:
    filepath = tmp_path / "test_doc.pdf"
    pdf_doc = fitz.open()
    page = pdf_doc.new_page(width=595, height=842)
    page.insert_text((72, 72), "PDF Test Document", fontsize=14)
    page.insert_text((72, 100), "This is a PDF with text content.", fontsize=11)
    pdf_doc.save(str(filepath))
    pdf_doc.close()
    return filepath


def _create_docx_file(tmp_path: Path) -> Path:
    filepath = tmp_path / "test_doc.docx"
    doc = DocxDocument()
    doc.add_heading("DOCX Test", level=1)
    doc.add_paragraph("This is a DOCX paragraph.")
    table = doc.add_table(rows=2, cols=2)
    table.style = "Table Grid"
    table.cell(0, 0).text = "A1"
    table.cell(0, 1).text = "B1"
    table.cell(1, 0).text = "A2"
    table.cell(1, 1).text = "B2"
    doc.save(str(filepath))
    return filepath


def _create_xlsx_file(tmp_path: Path) -> Path:
    filepath = tmp_path / "test_doc.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws["A1"] = "Name"
    ws["B1"] = "Value"
    ws["A2"] = "Foo"
    ws["B2"] = 42
    wb.save(str(filepath))
    wb.close()
    return filepath


def _create_pptx_file(tmp_path: Path) -> Path:
    filepath = tmp_path / "test_doc.pptx"
    try:
        from pptx import Presentation
    except ImportError:
        return filepath  # Will not create a valid file; tests should handle

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "PPTX Test"
    slide.placeholders[1].text = "Slide content."
    prs.save(str(filepath))
    return filepath


def _create_png_file(tmp_path: Path) -> Path:
    filepath = tmp_path / "test_doc.png"
    img = Image.new("RGB", (100, 50), color="white")
    img.save(str(filepath), format="PNG")
    return filepath


def _create_jpg_file(tmp_path: Path) -> Path:
    filepath = tmp_path / "test_doc.jpg"
    img = Image.new("RGB", (100, 50), color="white")
    img.save(str(filepath), format="JPEG")
    return filepath


# ============================================================
# Tests: parse() end-to-end for each format type
# ============================================================


class TestEndToEndParsing:
    def test_parse_txt(self, tmp_path: Path):
        filepath = _create_txt_file(tmp_path)
        result = parse(str(filepath))
        assert isinstance(result, ParseResult)
        assert len(result.blocks) >= 1
        assert "Hello World" in result.markdown

    def test_parse_md(self, tmp_path: Path):
        filepath = _create_md_file(tmp_path)
        result = parse(str(filepath))
        assert isinstance(result, ParseResult)
        assert "# Title" in result.markdown

    def test_parse_pdf(self, tmp_path: Path):
        filepath = _create_pdf_file(tmp_path)
        result = parse(str(filepath))
        assert isinstance(result, ParseResult)
        assert "PDF Test Document" in result.markdown

    def test_parse_docx(self, tmp_path: Path):
        filepath = _create_docx_file(tmp_path)
        result = parse(str(filepath))
        assert isinstance(result, ParseResult)
        assert "# DOCX Test" in result.markdown
        assert "A1" in result.markdown
        assert "B2" in result.markdown

    def test_parse_xlsx(self, tmp_path: Path):
        filepath = _create_xlsx_file(tmp_path)
        result = parse(str(filepath))
        assert isinstance(result, ParseResult)
        assert "Data" in result.markdown
        assert "Name" in result.markdown
        assert "Foo" in result.markdown

    def test_parse_pptx(self, tmp_path: Path):
        try:
            from pptx import Presentation  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")
        filepath = _create_pptx_file(tmp_path)
        result = parse(str(filepath))
        assert isinstance(result, ParseResult)
        assert "PPTX Test" in result.markdown
        assert "Slide 1" in result.markdown

    def test_parse_image_detection(self, tmp_path: Path):
        """Test image format detection and routing (not full OCR)."""
        filepath = _create_png_file(tmp_path)

        # Register a dummy OCR backend so OCR is available
        saved = dict(_BACKEND_REGISTRY)
        _BACKEND_REGISTRY.clear()
        try:
            register_backend("dummy_integration", _DummyOCRBackend)
            result = parse(str(filepath), ocr_backend="dummy_integration")
            assert isinstance(result, ParseResult)
            assert "Test OCR result" in result.markdown
        finally:
            _BACKEND_REGISTRY.clear()
            _BACKEND_REGISTRY.update(saved)


# ============================================================
# Tests: Error handling
# ============================================================


class TestErrorHandling:
    def test_unknown_format_raises_file_not_supported(self, tmp_path: Path):
        # Create a file with unknown extension
        filepath = tmp_path / "test.xyz"
        filepath.write_bytes(b"some content")

        with pytest.raises(FileNotSupportedError):
            parse(str(filepath))

    def test_nonexistent_file_raises(self):
        with pytest.raises(FileNotSupportedError):
            parse("/nonexistent/path/file.pdf")


# ============================================================
# Tests: FormatDetector accuracy
# ============================================================


class TestFormatDetectorAccuracy:
    def test_detect_pdf(self, tmp_path: Path):
        filepath = _create_pdf_file(tmp_path)
        result = detect_file_type(filepath)
        assert result == "pdf"

    def test_detect_docx(self, tmp_path: Path):
        filepath = _create_docx_file(tmp_path)
        result = detect_file_type(filepath)
        assert result == "docx"

    def test_detect_xlsx(self, tmp_path: Path):
        filepath = _create_xlsx_file(tmp_path)
        result = detect_file_type(filepath)
        assert result == "xlsx"

    def test_detect_pptx(self, tmp_path: Path):
        filepath = _create_pptx_file(tmp_path)
        if not filepath.exists():
            pytest.skip("pptx file not created")
        result = detect_file_type(filepath)
        assert result == "pptx"

    def test_detect_png(self, tmp_path: Path):
        filepath = _create_png_file(tmp_path)
        result = detect_file_type(filepath)
        assert result == "image"

    def test_detect_jpg(self, tmp_path: Path):
        filepath = _create_jpg_file(tmp_path)
        result = detect_file_type(filepath)
        assert result == "image"

    def test_detect_txt(self, tmp_path: Path):
        filepath = _create_txt_file(tmp_path)
        result = detect_file_type(filepath)
        assert result == "txt"

    def test_detect_md(self, tmp_path: Path):
        filepath = _create_md_file(tmp_path)
        result = detect_file_type(filepath)
        assert result == "md"

    def test_format_detector_coverage(self, tmp_path: Path):
        """Verify all 7 format types are detected correctly (accuracy check)."""
        formats = {
            "pdf": _create_pdf_file,
            "docx": _create_docx_file,
            "xlsx": _create_xlsx_file,
            "pptx": _create_pptx_file,
            "image": _create_png_file,
            "txt": _create_txt_file,
            "md": _create_md_file,
        }

        for expected_type, factory in formats.items():
            filepath = factory(tmp_path)
            if not filepath.exists():
                pytest.skip(f"File not created for {expected_type}")
            detected = detect_file_type(filepath)
            assert detected == expected_type, (
                f"Expected {expected_type} but got {detected} for {filepath.name}"
            )


# ============================================================
# Tests: CLI integration
# ============================================================


class TestCLIIntegration:
    def test_cli_txt_output(self, tmp_path: Path):
        """Test CLI parse with TXT file."""
        filepath = _create_txt_file(tmp_path)
        output_path = tmp_path / "output.md"

        import subprocess
        import sys

        result = subprocess.run(
            [sys.executable, "-m", "docforge.cli", "parse", str(filepath), "-o", str(output_path)],
            capture_output=True,
            text=True,
        )
        assert result.returncode == 0, f"CLI failed: {result.stderr}"
        assert output_path.exists()
        content = output_path.read_text(encoding="utf-8")
        assert "Hello World" in content

    def test_cli_md_output(self, tmp_path: Path):
        """Test CLI parse with MD file."""
        filepath = _create_md_file(tmp_path)
        output_path = tmp_path / "output_md.md"

        import subprocess
        import sys

        result = subprocess.run(
            [sys.executable, "-m", "docforge.cli", "parse", str(filepath), "-o", str(output_path)],
            capture_output=True,
            text=True,
        )
        assert result.returncode == 0, f"CLI failed: {result.stderr}"
        assert output_path.exists()
        content = output_path.read_text(encoding="utf-8")
        assert "Title" in content or "# Title" in content

    def test_cli_docx_output(self, tmp_path: Path):
        """Test CLI parse with DOCX file."""
        filepath = _create_docx_file(tmp_path)
        output_path = tmp_path / "output_docx.md"

        import subprocess
        import sys

        result = subprocess.run(
            [sys.executable, "-m", "docforge.cli", "parse", str(filepath), "-o", str(output_path)],
            capture_output=True,
            text=True,
        )
        assert result.returncode == 0, f"CLI failed: {result.stderr}"
        assert output_path.exists()
        content = output_path.read_text(encoding="utf-8")
        assert "DOCX Test" in content
